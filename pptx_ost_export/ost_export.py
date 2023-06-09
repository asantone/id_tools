# purpose: this tool merges extracts on-screen text (OST) from multiple .pptx files and stores the output in a .docx file
# author: Adam Santone, PhD
# date: 2023-06-09

#libraries
import glob, os                           # file utilities
from pptx_tools import *                  # conversion (slides to images)
from pptx import Presentation             # used to get the presenter notes
import docx                               # for creating MS Word .docx files
from docx.enum.section import WD_ORIENT   # for Word document layout
from docx.enum.section import WD_SECTION  # for Word document layout
from docx.shared import Cm                # metric system
from docx.shared import RGBColor          # for hexadecimal color notation
#import re                                 # regular expressions for XML compatibility line
from natsort import natsorted             # for natural sorting order (1,2,3...n)

#generate natural-sorted list of files to be processed
pptList = []
pptList = natsorted(glob.glob('./**/*.pptx', recursive=True))

#on-screen text export for proofreading
def ost(pptx_name):

    ppt = Presentation(pptx_name)

    #create a Word Document (with optional header)
    document = docx.Document()
    #document.add_heading('On-screen Text (OST)', level=1)
    #document.add_heading('Course Information', level=3)
    #document.add_paragraph('Course Code: ', style='List Bullet')
    #document.add_paragraph('Course Title: ', style='List Bullet')
    #document.add_heading('Module Information', level=3)
    #document.add_paragraph('Module Number: ', style='List Bullet')
    #document.add_paragraph('Module Title: ', style='List Bullet')

    for slide in ppt.slides:
        idx = ppt.slides.index(slide)+1
        document.add_heading("slide_" + str(idx), level=3) #add the slide number to a header
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                try:
                    document.add_paragraph(shape.text)  #get the text out of the text box
                except:
                    pass
            else:
                pass

    #save Word document
    fn = os.path.splitext(os.path.basename(i))[0]
    filename = "ost_" + fn + ".docx"
    document.save(filename)
    print ("Complete: ", filename)
    print ("~ ~ ~ ~ ~")


for i in pptList:
    ost(i)
