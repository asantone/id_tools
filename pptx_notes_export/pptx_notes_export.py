# purpose: this tool exports .pptx presenter notes (below slides) to .docx 
# author: Adam Santone, PhD
# date: 2023-06-09

#libraries
import glob, os                           # file utilities
from pptx_tools import *                  # conversion (slides to images)
from pptx import Presentation             # used to get the presenter notes
import docx                               # for creating MS Word .docx files
from docx.enum.section import WD_ORIENT   # for Word document layout
from docx.enum.section import WD_SECTION  # for Word document layout
from docx.shared import Cm                # metric system
from docx.shared import RGBColor          # for hexadecimal color notation
import re                                 # regular expressions for XML compatibility line
from natsort import natsorted             # for natural sorting order (1,2,3...n)
from PIL import Image                     # image manipulation


#generate natural-sorted list of files to be processed
pptList = natsorted(glob.glob('./**/*.pptx', recursive=True))

#empty array to hold the notes
slideInfo = {}

# #function to export scripts
def scripts(pptx_name):

    #create a Word Document
    document = docx.Document()

    #add the header details (optional)
    document.add_heading('Script', level=1)

    #document.add_heading('Course Information', level=3)
    #document.add_paragraph('Course Code: ', style='List Bullet')
    #document.add_paragraph('Course Title: ', style='List Bullet')

    #document.add_heading('Module Information', level=3)
    #document.add_paragraph('Module Number: ', style='List Bullet')
    #document.add_paragraph('Module Title: ', style='List Bullet')

    #document.add_heading('Filming Information', level=3)
    #document.add_paragraph('Filming Date (YYYY-MM-DD): ', style='List Bullet')
    #document.add_paragraph('Additional Notes: ', style='List Bullet')

    #loop through image list and add them to the Word document
    for i, t in enumerate(slideInfo):
        document.add_heading("@_" + out + "_slide_" + list(slideInfo.keys())[i], level=3)
        document.add_paragraph(re.sub(u'[^\u0020-\uD7FF\u0009\u000A\u000D\uE000-\uFFFD\U00010000-\U0010FFFF]+', '', list(slideInfo.values())[i]))  #force XML compatibility

    #save Word document
    filename = "script_" + pptx_name + ".docx"
    document.save(filename)
    print ("Complete: ", filename)
    #print ("~ ~ ~ ~ ~")


#loop through pptx files
for i in pptList:
    #SETUP
    out = os.path.splitext(os.path.basename(i))[0]   #truncates the path to just the filename; also used to identify the current file in the two main functions
    path = os.getcwd()
    prePath = path.replace('\\','\\\\') #use double slashes
    postPath = ".pptx"                               #forced path stuff
    inFile = (prePath + "\\\\" + out + postPath)     #create the full path for the input .pptx file
    outFolder = (prePath + "\\\\" + out)             #create the full path for the output folder (same as input filename)

    #CONVERT SLIDES TO NOTES
    ppt=Presentation(i)
    for page, slide in enumerate(ppt.slides):
        textNote = slide.notes_slide.notes_text_frame.text
        slideInfo[str(page)] = textNote

    #generate output
    scripts(out)                    

